


import glob
import re
import numpy as np
import yaml
import textwrap
import shutil
import os
import io
import pickle
import matplotlib.pyplot as plt
from collections import defaultdict
from PIL import Image, ImageDraw, ImageFont
from typing import Optional
import importlib
import wandb
from enum import Enum
import argparse


######################## PPTX ########################
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER

ppt = Presentation("test-shift-final/results/test-shift-template.pptx")

# sorted x, y
def get_placeholders(slide):
    # Initialize lists for image and text placeholders
    image_placeholders = []
    text_placeholders = []

    # Iterate through shapes in the slide
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type

            # Check placeholder type
            if placeholder_type == PP_PLACEHOLDER.PICTURE:  # PP_PLACEHOLDER.PICTURE (6) or PP_PLACEHOLDER.CONTENT (7)
                image_placeholders.append(shape)
            elif placeholder_type == PP_PLACEHOLDER.BODY:  # PP_PLACEHOLDER.TITLE (1), SUBTITLE (2), TEXT (3, 4)
                text_placeholders.append(shape)
    image_placeholders.sort(key=lambda x: (x.top, x.left))
    text_placeholders.sort(key=lambda x: (x.top, x.left))

    return image_placeholders, text_placeholders

def fill_contents(slide, images, texts):
    image_placeholders, text_placeholders = get_placeholders(slide)
    assert len(image_placeholders) == len(images)
    assert len(text_placeholders) == len(texts)
    for image_placeholder, image in zip(image_placeholders, images):
        if image is not None:
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG")
            image_placeholder.insert_picture(buffer)
    for text_placeholder, text in zip(text_placeholders, texts):
        text_placeholder.text = text

######################## PPTX ########################

################### different mode ###################
path_list = [
    "test-shift-final/results/mode=1/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=2/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=3/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=1/p=True/std=1/{filename}.jpeg",
    "test-shift-final/results/mode=1/p=True/std=3/{filename}.jpeg",
]

slide_images = []
slide_texts = []

for path in path_list:
    target_path = path.format(filename=f"target_images")
    slide_images.append(Image.open(target_path))
    for index in range(0,101,10):
        images_path = path.format(filename=f"images_{index}")
        slide_images.append(Image.open(images_path))

layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(layout)
fill_contents(slide, slide_images, slide_texts)

################### different mode ###################

################# different mode-1-k #################
path_list = [
    "test-shift-final/results/mode=1-k=8/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=1-k=4/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=1-k=2/p=True/std=0/{filename}.jpeg",
    "test-shift-final/results/mode=1-k=1/p=True/std=0/{filename}.jpeg"
]

slide_images = []
slide_texts = []

for path in path_list:
    target_path = path.format(filename=f"target_images")
    slide_images.append(Image.open(target_path))
    for index in range(0,101,10):
        images_path = path.format(filename=f"images_{index}")
        slide_images.append(Image.open(images_path))

layout = ppt.slide_layouts[1]
slide = ppt.slides.add_slide(layout)
fill_contents(slide, slide_images, slide_texts)
################# different mode-1-k #################


ppt.save(f"test-shift-final/results/test-shift.pptx")