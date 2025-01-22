


import glob
import re
import numpy as np
import yaml
import textwrap
import shutil
import os
import io
import pickle
import matplotlib.pyplot as plt
from collections import defaultdict
from PIL import Image, ImageDraw, ImageFont
from typing import Optional
import importlib
import wandb
from enum import Enum
import argparse

dir_path = os.path.dirname(os.path.realpath(__file__))

# Create a parser object
parser = argparse.ArgumentParser(description="A script to demonstrate taking arguments.")
parser.add_argument("--demo-id", type=int, default=0, help="Demo ID.")
args = parser.parse_args()

def cache_key(algo,name):
    return f"{algo}-{name}"

all_images = defaultdict(list)

name_to_prompts_map = {
    "deerelephant": "A yellow reindeer and a blue elephant.",
    "trafficlight": "A traffic light with yellow at top, green at middle, and red at bottom.",
    "apple": "Seven red apples arranged in a circle.",
    "cyberdog": "A natural fluffy dog talking to a cybertic dog.",
    "puppynose": "Side view of a puppy lying on floor, one butterfly stopping on its nose.",
    "robotplant": "A cute cybernetic robot plants a tree in the forest.",
    "ocean": "A helicopter floating under the ocean.",
    "sandglass": "A transparent glass filled with a mixture of water and sand, with one feather floating inside.",
    "penguin": "A photo realistic photo showing a penguin standing on a very small floating ice, with a tree is on fire in the background.",
    "basket": "Exactly one orange in a basket of apples.",
    "icecube": "A glass of water with exactly one ice cube.",
    "catbutterfly": "A cat with butterfly wings.",
}

names = [
    "deerelephant", "trafficlight", "apple", "cyberdog", "puppynose", "robotplant", "ocean", "sandglass", "penguin", "basket", "icecube", "catbutterfly",
    "compress", "incompress", "aesthetic",
]

sr_names = ["compress", "incompress", "aesthetic"]
prompt_names = list(set(names) - set(sr_names))

base_demo_id = args.demo_id

name_demo_ids_map = {
    "deerelephant" : [2],
    "trafficlight" : [0],
    "apple"        : [10],
    "cyberdog"     : [10],
    "puppynose"    : [10],
    "robotplant"   : [10],
    "ocean"        : [0],
    "sandglass"    : [6],
    "penguin"      : [0],
    "basket"       : [10],
    "icecube"      : [5],
    "catbutterfly" : [10],
    "compress"     : [0,1,2,3,4,5,6,7,8,9,10,11,12,13,14,15],
    "incompress"   : [0,1,2,3,4,5,6,7,8,9,10,11,12,13,14,15],
    "aesthetic"    : [0,1,2,3,4,5,6,7,8,9,10,11,12,13,14,15],
}

# 0: eval until 200
# 1: eval until 50
class LAYOUT_NAME(Enum):
    LONG    = 0
    SHORT   = 1
    SR      = 2
    SR_20  = 3
name_slide_layout_map = {
    "deerelephant": LAYOUT_NAME.LONG,
    "trafficlight": LAYOUT_NAME.LONG,
    "apple"       : LAYOUT_NAME.LONG,
    "cyberdog"    : LAYOUT_NAME.SHORT,
    "puppynose"   : LAYOUT_NAME.SHORT,
    "robotplant"  : LAYOUT_NAME.SHORT,
    "ocean"       : LAYOUT_NAME.SHORT,
    "sandglass"   : LAYOUT_NAME.SHORT,
    "penguin"     : LAYOUT_NAME.SHORT,
    "basket"      : LAYOUT_NAME.SHORT,
    "icecube"     : LAYOUT_NAME.SHORT,
    "catbutterfly": LAYOUT_NAME.SHORT,
    "compress"    : LAYOUT_NAME.SR,
    "incompress"  : LAYOUT_NAME.SR,
    "aesthetic"   : LAYOUT_NAME.SR_20,
}

import argparse

parser = argparse.ArgumentParser()

cache_path = f"{dir_path}/wandb_cache.pkl"
with open(cache_path, 'rb') as f:
    wandb_cache = pickle.load(f)

def random_image():
    return Image.fromarray(np.random.randint(0, 256, (1024, 1024, 3), dtype=np.uint8))

api = wandb.Api()
wandb_path = "kimyong95/guide-stable-diffusion"

missing = []

######################## PPTX ########################
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER

ppt = Presentation(f"{dir_path}/visualize/visualize-template.pptx")

# sorted x, y
def get_placeholders(slide):
    # Initialize lists for image and text placeholders
    image_placeholders = []
    text_placeholders = []

    # Iterate through shapes in the slide
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type

            # Check placeholder type
            if placeholder_type == PP_PLACEHOLDER.PICTURE:  # PP_PLACEHOLDER.PICTURE (6) or PP_PLACEHOLDER.CONTENT (7)
                image_placeholders.append(shape)
            elif placeholder_type == PP_PLACEHOLDER.BODY:  # PP_PLACEHOLDER.TITLE (1), SUBTITLE (2), TEXT (3, 4)
                text_placeholders.append(shape)
    image_placeholders.sort(key=lambda x: (x.top, x.left))
    text_placeholders.sort(key=lambda x: (x.top, x.left))

    return image_placeholders, text_placeholders

def fill_contents(slide, images, texts):
    image_placeholders, text_placeholders = get_placeholders(slide)
    assert len(image_placeholders) == len(images)
    assert len(text_placeholders) == len(texts)
    for image_placeholder, image in zip(image_placeholders, images):
        if image is not None:
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG")
            image_placeholder.insert_picture(buffer)
    for text_placeholder, text in zip(text_placeholders, texts):
        if type(text) == str:
            text_placeholder.text = text
        elif type(text) == list:
            text_frame = text_placeholder.text_frame
            for i, _text in enumerate(text):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = _text["text"]
                p.font.size = Pt(_text["fontsize"])

######################## PPTX ########################


######################## OURS ########################

# i: column
# j: row
def crop_image(image, i,j, grid_size=1024):
    width, height = image.size
    assert width == height
    start_i = i * grid_size + (i+1)*2
    start_j = j * grid_size + (j+1)*2
    image = image.crop((start_i, start_j, start_i + grid_size, start_j + grid_size))
    return image

def concat_images(image_matrix):

    # assert 2D matrix
    assert all(len(row) == len(image_matrix[0]) for row in image_matrix)

    # Concatenate images within each row horizontally
    rows = []
    for image_row in image_matrix:
        widths, heights = zip(*(i.size for i in image_row))
        total_width = sum(widths)
        max_height = max(heights)

        new_row = Image.new('RGB', (total_width, max_height))
        x_offset = 0
        for img in image_row:
            new_row.paste(img, (x_offset, 0))
            x_offset += img.size[0]
        rows.append(new_row)

    # Concatenate rows vertically
    widths, heights = zip(*(i.size for i in rows))
    max_width = max(widths)
    total_height = sum(heights)

    final_concat = Image.new('RGB', (max_width, total_height))
    y_offset = 0
    for row in rows:
        final_concat.paste(row, (0, y_offset))
        y_offset += row.size[1]

    return final_concat

def get_validation_image(run_path, epoch):
    image_path = run_path + f"/ablation/validation/{epoch}/l={epoch}.jpg"
    if not os.path.exists(image_path):
        missing.append(image_path)
        return random_image()
    image = Image.open(image_path)
    return image

def get_validation_texts(run_path, epoch):
    text_path = run_path + f"/ablation/validation/{epoch}/llava.txt"
    with open(text_path, "r") as f:
        texts = f.read().splitlines()
    return texts

def get_prompt(run_path):
    config_path = run_path + "/files/config-cli.yaml"
    with open(config_path, "r") as f:
        config = yaml.safe_load(f)
    return config["model"]["init_args"]["generate_prompt"]

algos = ["ours", "ours-ddim"]

def copy_compress(src, dest, dest_compress):
    if not os.path.exists(src):
        missing.append(src)
        return
    
    if not os.path.exists(dest):
        os.makedirs(os.path.dirname(dest), exist_ok=True)
        shutil.copy(src, dest)
    
    if not os.path.exists(dest_compress):
        os.makedirs(os.path.dirname(dest_compress), exist_ok=True)
        Image.open(src).save(
            dest_compress,
            quality=22,
            optimize=True
        )

for algo in algos:
    for index, name in enumerate(names):
        if name_slide_layout_map[name] == LAYOUT_NAME.SR_20:
            includes_epoches = [0,19,39,59,79,99]
        else:
            includes_epoches = [0,9,19,29,39,49]
        
        _cache_key = cache_key(algo,name)
        _run_id = wandb_cache[_cache_key]["run_id"]
        demo_ids = name_demo_ids_map[name]

        paths = glob.glob(f"wandb/run-*-{_run_id}")
        assert len(paths) <= 1, f"Multiple runs found: {paths}"
        path = paths[0] if paths else None

        # save first steps
        copy_compress(
            src=f"{path}/ablation/train/{includes_epoches[-1]}/l=0.jpg",
            dest=f"{dir_path}/{algo}-init/{index}.jpeg",
            dest_compress=f"{dir_path}/{algo}-init-compress/{index}.jpeg"
        )

        # save final steps
        copy_compress(
            src=f"{path}/ablation/train/{includes_epoches[-1]}/l={includes_epoches[-1]}.jpg",
            dest=f"{dir_path}/{algo}-final/{index}.jpeg",
            dest_compress=f"{dir_path}/{algo}-final-compress/{index}.jpeg"
        )

        # copy validation
        if name in sr_names:
            copy_compress(
                src=f"{path}/ablation/validation/{includes_epoches[-1]}/l=0.jpg",
                dest=f"{dir_path}/{algo}-init/{index}-val.jpeg",
                dest_compress=f"{dir_path}/{algo}-init-compress/{index}-val.jpeg"
            )
            copy_compress(
                src=f"{path}/ablation/validation/{includes_epoches[-1]}/l={includes_epoches[-1]}.jpg",
                dest=f"{dir_path}/{algo}-final/{index}-val.jpeg",
                dest_compress=f"{dir_path}/{algo}-final-compress/{index}-val.jpeg"
            )        
        
        for _demo_id in demo_ids:

            paths = glob.glob(f"wandb/run-*-{_run_id}")
            assert len(paths) <= 1, f"Multiple runs found: {paths}"
            path = paths[0] if paths else None

            image_row = []
            for e in includes_epoches:
                if path is not None:
                    image = get_validation_image(path, e)
                    demo_row = _demo_id // 4
                    demo_col = _demo_id % 4
                    image = crop_image(image, demo_col, demo_row, grid_size=1024)
                else:
                    missing.append(f"wandb/run-*-{_run_id}/ablation/validation/{e}/l={e}.jpg")
                    image = random_image()
                image_row.append(image)
            all_images[_cache_key].append(image_row)

######################## OURS ########################

######################## DNO #########################
dno_names = list(set(names) - set(sr_names))
for name in dno_names:
    demo_ids = name_demo_ids_map[name]
    for _demo_id in demo_ids:
        _name = f"{name}-index={_demo_id}"
        _cache_key = cache_key("dno",_name)
        _run_id = wandb_cache[_cache_key]["run_id"]

        paths = glob.glob(f"related_works/Direct-Noise-Optimization/wandb/run-*-{_run_id}")
        assert len(paths) <= 1, f"Multiple runs found: {paths}"
        path = paths[0] if paths else None

        images_list = []
        if name_slide_layout_map[name] == LAYOUT_NAME.LONG:
            epoches = [0,9,19,29,39,49,499]
        elif name_slide_layout_map[name] == LAYOUT_NAME.SHORT:
            epoches = [0,9,19,29,39,49]

        for e in epoches:
            image_path = f"{path}/images/validation/{e}.jpg"
            if path is not None and os.path.exists(image_path):
                image = Image.open(image_path)
            else:
                missing.append(f"related_works/Direct-Noise-Optimization/wandb/run-*-{_run_id}/images/validation/{e}.jpg")
                image = random_image()
            images_list.append(image)

        all_images[_cache_key].append(images_list)

######################## DNO #########################

##################### baselines ######################
algos = ["ddpo", "dpok", "d3po"]

for algo in algos:
    for name in names:
        _cache_key = cache_key(algo,name)
        demo_ids = name_demo_ids_map[name]
        for _demo_id in demo_ids:

            # load pre-train image
            images_list = [Image.open(f"{dir_path}/ppo-pretrain-images/{name}/{_demo_id}.jpeg")]
            
            image_ = wandb_cache[_cache_key]["history"]["validation/images"]
            image_k = wandb_cache[_cache_key]["history"][image_.notnull()]["epoch"].values.tolist()
            image_v = image_[image_.notnull()].values.tolist()
            image_dict = dict(zip(image_k, image_v))

            image_dir = f"{dir_path}/wandb-images/{_cache_key}"

            if name_slide_layout_map[name] in [LAYOUT_NAME.LONG, LAYOUT_NAME.SR]:
                epoches = [9,19,29,39,49,499]
            elif name_slide_layout_map[name] == LAYOUT_NAME.SHORT:
                epoches = [9,19,29,39,49]
            elif name_slide_layout_map[name] == LAYOUT_NAME.SR_20:
                epoches = [19,39,69,79,99,499]
            
            _wandb_run = None
            for e in epoches:
                image = None
                if e in image_dict:
                    image_name = image_dict[e]["filenames"][_demo_id]
                    image_path = f"{image_dir}/{image_name}"
                    if os.path.exists(image_path):
                        image = Image.open(image_path)
                    else:
                        print(f"Downloading files [{_cache_key}] ...")
                        if _wandb_run is None:
                            _run_id = wandb_cache[_cache_key]["run_id"]
                            _wandb_run = api.run(f"{wandb_path}/{_run_id}")
                        _wandb_run.file(name=image_name).download(image_dir)
                        image = Image.open(image_path)

                images_list.append(image if image else random_image())

            all_images[_cache_key].append(images_list)

##################### baselines ######################

with open(f'{dir_path}/visualize/missing.txt', 'a') as f:
    for item in missing:
        f.write(f"{item}\n")

################### create slides ####################
algos = ["ddpo", "dpok", "d3po", "dno", "ours-ddim", "ours"]
algo_label_map = {
    "ddpo": "DDPO",
    "dpok": "DPOK",
    "d3po": "D3PO",
    "dno": "DNO",
    "ours-ddim": [
        { "text": "Fast Direct", "fontsize": 36 },
        { "text": "(w/ DDIM)", "fontsize": 20},
    ],
    "ours": [
        { "text": "Fast Direct", "fontsize": 36 },
        { "text": "(w/ EDM)", "fontsize": 20},
    ]
}

three_dots = Image.open(f"{dir_path}/assets/three-dots.jpg")
for name in names:
    demo_ids = name_demo_ids_map[name]
    for i, _demo_id in enumerate(demo_ids):
        layout = ppt.slide_layouts[name_slide_layout_map[name].value]
        slide = ppt.slides.add_slide(layout)
        slide_images = []
        slide_texts = []
        for algo in algos:
            if algo == "dno" and name in sr_names:
                continue
            if algo == "dno":
                cache_key_ = cache_key("dno",f"{name}-index={_demo_id}")
            else:
                cache_key_ = cache_key(algo,name)
            assert len(all_images[cache_key_]) == len(demo_ids)
            images_ = all_images[cache_key_][i]

            if name_slide_layout_map[name] in [LAYOUT_NAME.LONG, LAYOUT_NAME.SR, LAYOUT_NAME.SR_20]:
                if algo in ["ddpo", "dpok", "d3po", "dno"]:
                    slide_images_row = images_[:-1] + [three_dots] + [images_[-1]]
                elif algo in ["ours", "ours-ddim"]:
                    slide_images_row = images_ + [None,None]
            else:
                slide_images_row = images_

            slide_images.extend(slide_images_row)
            slide_texts.extend([algo_label_map[algo]])

        fill_contents(slide, slide_images, slide_texts)
ppt.save(f"{dir_path}/visualize/visualize.pptx")

################### create slides ####################